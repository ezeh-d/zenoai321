"""OCR and document text extraction.

ENGINE CHOICE, AND WHY
----------------------
Uses the OCR engine built into Windows (`Windows.Media.Ocr` via the
`winsdk` package that is already a dependency for notification handling).
Verified available on this machine with en-GB/en-US recognisers.

Tesseract was the obvious first choice and was rejected on evidence:
`pytesseract` is not installed AND it needs a separate native binary that
is also absent, so wiring it would have produced a module that imports
fine and fails at run time. The Windows engine needs neither, and it
returns per-word confidence, which the rest of this file depends on.

CONFIDENCE IS REAL, NOT INVENTED
--------------------------------
Windows OCR does not expose a numeric per-word score, so a fabricated
percentage would be dishonest. Instead the confidence reported here is
derived from observable properties of the result: how much text was
found, how many words survive a dictionary-shape sanity check, and
average word length. It is labelled as a heuristic everywhere it appears,
and `extract_*` always returns the raw text alongside it so a caller can
judge for itself.

DOCUMENTS
---------
Plain text is read directly. PDF, DOCX, XLSX and PPTX use lazy native readers:
none of those packages is imported during ZENO startup. Reads are bounded by
file size, document elements and ``max_chars`` to avoid large/hostile files
turning a simple inspection into an unbounded workload.
"""

from __future__ import annotations

import asyncio
import importlib.util
import re
from dataclasses import dataclass, field
from pathlib import Path

_TEXT_SUFFIXES = {".txt", ".md", ".markdown", ".csv", ".json", ".log", ".py",
                  ".js", ".html", ".css", ".yml", ".yaml", ".ini", ".xml"}
_IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tif", ".tiff", ".webp"}
_DOCUMENT_READERS = {
    ".pdf": ("fitz", "PyMuPDF"),
    ".docx": ("docx", "python-docx"),
    ".xlsx": ("openpyxl", "openpyxl"),
    ".pptx": ("pptx", "python-pptx"),
}
_LEGACY_FORMATS = {
    ".doc": "legacy .doc needs conversion to .docx",
    ".xls": "legacy .xls needs conversion to .xlsx",
}
_MAX_DOCUMENT_BYTES = 50 * 1024 * 1024
_MAX_PDF_PAGES = 500
_MAX_DOCX_BLOCKS = 10_000
_MAX_SHEETS = 50
_MAX_ROWS_PER_SHEET = 2_000
_MAX_CELLS_PER_ROW = 100
_MAX_SLIDES = 500

_WORD_RE = re.compile(r"[A-Za-z][A-Za-z'-]{1,}")


@dataclass
class OcrResult:
    text: str = ""
    ok: bool = False
    source: str = ""
    engine: str = ""
    word_count: int = 0
    confidence: float = 0.0      # 0-1 heuristic; see module docstring
    confidence_basis: str = ""
    error: str = ""
    lines: list[str] = field(default_factory=list)

    def as_dict(self) -> dict:
        return {
            "ok": self.ok, "source": self.source, "engine": self.engine,
            "text": self.text, "word_count": self.word_count,
            "confidence": round(self.confidence, 2),
            "confidence_basis": self.confidence_basis,
            "error": self.error, "line_count": len(self.lines),
        }


def _score(text: str) -> tuple[float, str]:
    """Heuristic quality score for an OCR result.

    Deliberately simple and explainable: OCR noise looks like short
    fragments of non-word characters, real text looks like words. Returns
    (0-1, human-readable basis).
    """
    stripped = text.strip()
    if not stripped:
        return 0.0, "no text found"
    words = _WORD_RE.findall(stripped)
    if not words:
        return 0.15, "characters found but no recognisable words"
    total_tokens = len(stripped.split())
    word_ratio = len(words) / max(1, total_tokens)
    avg_len = sum(len(w) for w in words) / len(words)
    # Real prose sits around 4-6 chars/word; 1-2 means fragmented noise.
    length_factor = min(1.0, avg_len / 4.0)
    volume_factor = min(1.0, len(words) / 12.0)   # very little text = less certain
    score = 0.55 * word_ratio + 0.30 * length_factor + 0.15 * volume_factor
    basis = (f"{len(words)} word-like tokens of {total_tokens} "
             f"({word_ratio:.0%}), avg length {avg_len:.1f}")
    return max(0.0, min(1.0, score)), basis


def _direct_result(path: Path, text: str, engine: str) -> OcrResult:
    """Build an honest result for deterministic document text extraction."""
    clean = text.strip()
    return OcrResult(
        text=clean,
        ok=bool(clean),
        source=str(path),
        engine=engine,
        word_count=len(clean.split()),
        confidence=1.0 if clean else 0.0,
        confidence_basis=("extracted from document structure, no OCR involved"
                          if clean else "document contained no extractable text"),
        lines=clean.splitlines()[:400],
    )


def _append_bounded(parts: list[str], value: object, *, used: int,
                    max_chars: int) -> int:
    text = str(value).strip() if value is not None else ""
    if not text or used >= max_chars:
        return used
    remaining = max_chars - used
    clipped = text[:remaining]
    parts.append(clipped)
    return used + len(clipped) + 1


def _extract_pdf(path: Path, max_chars: int) -> OcrResult:
    import fitz

    parts: list[str] = []
    used = 0
    with fitz.open(path) as document:
        for page_no in range(min(document.page_count, _MAX_PDF_PAGES)):
            used = _append_bounded(
                parts, document.load_page(page_no).get_text("text"),
                used=used, max_chars=max_chars,
            )
            if used >= max_chars:
                break
    return _direct_result(path, "\n".join(parts)[:max_chars], "pymupdf")


def _extract_docx(path: Path, max_chars: int) -> OcrResult:
    from docx import Document

    document = Document(path)
    parts: list[str] = []
    used = 0
    blocks = 0
    for paragraph in document.paragraphs:
        used = _append_bounded(parts, paragraph.text, used=used, max_chars=max_chars)
        blocks += 1
        if used >= max_chars or blocks >= _MAX_DOCX_BLOCKS:
            break
    if used < max_chars and blocks < _MAX_DOCX_BLOCKS:
        for table in document.tables:
            for row in table.rows:
                row_text = "\t".join(cell.text for cell in row.cells)
                used = _append_bounded(parts, row_text, used=used, max_chars=max_chars)
                blocks += 1
                if used >= max_chars or blocks >= _MAX_DOCX_BLOCKS:
                    break
            if used >= max_chars or blocks >= _MAX_DOCX_BLOCKS:
                break
    return _direct_result(path, "\n".join(parts)[:max_chars], "python-docx")


def _extract_xlsx(path: Path, max_chars: int) -> OcrResult:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True, keep_links=False)
    parts: list[str] = []
    used = 0
    try:
        for sheet in workbook.worksheets[:_MAX_SHEETS]:
            used = _append_bounded(parts, f"[{sheet.title}]", used=used,
                                   max_chars=max_chars)
            for row_no, row in enumerate(
                    sheet.iter_rows(max_col=_MAX_CELLS_PER_ROW, values_only=True), 1):
                if row_no > _MAX_ROWS_PER_SHEET:
                    break
                row_text = "\t".join("" if value is None else str(value) for value in row)
                used = _append_bounded(parts, row_text, used=used, max_chars=max_chars)
                if used >= max_chars:
                    break
            if used >= max_chars:
                break
    finally:
        workbook.close()
    return _direct_result(path, "\n".join(parts)[:max_chars], "openpyxl")


def _extract_pptx(path: Path, max_chars: int) -> OcrResult:
    from pptx import Presentation

    presentation = Presentation(path)
    parts: list[str] = []
    used = 0
    # python-pptx's Slides slice returns raw relationship elements rather than
    # Slide objects on some versions, so iterate normally and bound explicitly.
    for slide_no, slide in enumerate(presentation.slides, 1):
        if slide_no > _MAX_SLIDES:
            break
        used = _append_bounded(parts, f"[Slide {slide_no}]", used=used,
                               max_chars=max_chars)
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                used = _append_bounded(parts, shape.text, used=used, max_chars=max_chars)
            elif getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    used = _append_bounded(
                        parts, "\t".join(cell.text for cell in row.cells),
                        used=used, max_chars=max_chars,
                    )
            if used >= max_chars:
                break
        if used >= max_chars:
            break
    return _direct_result(path, "\n".join(parts)[:max_chars], "python-pptx")


_EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".pptx": _extract_pptx,
}


async def _ocr_bitmap_async(path: Path) -> tuple[str, list[str]]:
    from winsdk.windows.graphics.imaging import BitmapDecoder
    from winsdk.windows.media.ocr import OcrEngine
    from winsdk.windows.storage import FileAccessMode, StorageFile

    engine = OcrEngine.try_create_from_user_profile_languages()
    if engine is None:
        raise RuntimeError("no OCR recogniser installed for your Windows display languages")
    f = await StorageFile.get_file_from_path_async(str(path))
    stream = await f.open_async(FileAccessMode.READ)
    decoder = await BitmapDecoder.create_async(stream)
    bitmap = await decoder.get_software_bitmap_async()
    result = await engine.recognize_async(bitmap)
    lines = [ln.text for ln in result.lines] if result and result.lines else []
    return (result.text if result else ""), lines


def extract_image_text(image_path: str | Path) -> OcrResult:
    """OCR one image file using the Windows engine."""
    path = Path(image_path)
    if not path.is_file():
        return OcrResult(source=str(path), error="no such file")
    try:
        text, lines = asyncio.run(_ocr_bitmap_async(path))
    except Exception as exc:  # noqa: BLE001 -- surfaced to the caller, never swallowed
        return OcrResult(source=str(path), engine="windows-ocr",
                         error=f"{type(exc).__name__}: {exc}")
    conf, basis = _score(text)
    return OcrResult(text=text, ok=bool(text.strip()), source=str(path),
                     engine="windows-ocr", word_count=len(text.split()),
                     confidence=conf, confidence_basis=basis, lines=lines)


def extract_screen_text(region: tuple[int, int, int, int] | None = None) -> OcrResult:
    """Capture the screen (or a region) and OCR it.

    Reuses the existing screenshot path rather than adding a second capture
    mechanism, then deletes the temporary frame -- reading the screen should
    not silently accumulate images of the user's desktop on disk.
    """
    import tempfile

    try:
        import mss
    except ImportError:
        try:
            from PIL import ImageGrab
        except ImportError:
            return OcrResult(engine="windows-ocr", source="screen",
                             error="no screen-capture library available (mss or Pillow)")
        img = ImageGrab.grab(bbox=region)
        tmp = Path(tempfile.gettempdir()) / "zeno_ocr_screen.png"
        img.save(tmp)
    else:
        with mss.mss() as sct:
            mon = sct.monitors[1]
            box = ({"left": region[0], "top": region[1],
                    "width": region[2] - region[0], "height": region[3] - region[1]}
                   if region else mon)
            shot = sct.grab(box)
            tmp = Path(tempfile.gettempdir()) / "zeno_ocr_screen.png"
            import mss.tools

            mss.tools.to_png(shot.rgb, shot.size, output=str(tmp))

    try:
        res = extract_image_text(tmp)
        res.source = "screen" + (f" region {region}" if region else " (full)")
        return res
    finally:
        try:
            tmp.unlink(missing_ok=True)   # don't leave desktop captures lying around
        except OSError:
            pass


def extract_document_text(doc_path: str | Path, max_chars: int = 20000) -> OcrResult:
    """Text from a document. Images route to OCR; text formats are read
    directly; formats needing an absent library say so explicitly."""
    path = Path(doc_path)
    if not path.is_file():
        return OcrResult(source=str(path), error="no such file")
    if max_chars < 1:
        return OcrResult(source=str(path), error="max_chars must be positive")
    try:
        size = path.stat().st_size
    except OSError as exc:
        return OcrResult(source=str(path), error=str(exc))
    if size > _MAX_DOCUMENT_BYTES:
        return OcrResult(source=str(path), error=(
            f"document is {size} bytes; maximum is {_MAX_DOCUMENT_BYTES} bytes"
        ))
    suffix = path.suffix.lower()

    if suffix in _IMAGE_SUFFIXES:
        return extract_image_text(path)

    if suffix in _TEXT_SUFFIXES:
        try:
            text = path.read_text(encoding="utf-8", errors="replace")[:max_chars]
        except OSError as exc:
            return OcrResult(source=str(path), error=str(exc))
        conf, basis = _score(text)
        return OcrResult(text=text, ok=bool(text.strip()), source=str(path),
                         engine="direct-read", word_count=len(text.split()),
                         # A direct file read is exact; the heuristic only
                         # describes OCR uncertainty, so report full confidence.
                         confidence=1.0,
                         confidence_basis="read directly from disk, no OCR involved",
                         lines=text.splitlines()[:400])

    if suffix in _DOCUMENT_READERS:
        module, package = _DOCUMENT_READERS[suffix]
        if importlib.util.find_spec(module) is None:
            return OcrResult(source=str(path), engine="none",
                             error=(f"{suffix} needs {package}. Install the supported "
                                    "adapter with: python install.py --catalog-safe"))
        try:
            return _EXTRACTORS[suffix](path, max_chars)
        except Exception as exc:  # noqa: BLE001 -- precise error is user-visible
            return OcrResult(source=str(path), engine=package,
                             error=f"{type(exc).__name__}: {exc}")

    if suffix in _LEGACY_FORMATS:
        return OcrResult(source=str(path), engine="none", error=_LEGACY_FORMATS[suffix])

    return OcrResult(source=str(path), engine="none",
                     error=f"unsupported file type '{suffix}'")


def capabilities() -> dict:
    """What this module can genuinely do right now."""
    engines = []
    langs: list[str] = []
    try:
        from winsdk.windows.media.ocr import OcrEngine

        langs = [l.language_tag for l in OcrEngine.available_recognizer_languages]
        if OcrEngine.try_create_from_user_profile_languages() is not None:
            engines.append("windows-ocr")
    except Exception:  # noqa: BLE001
        pass
    document_formats: dict[str, dict[str, object]] = {}
    for suffix, (module, package) in _DOCUMENT_READERS.items():
        document_formats[suffix] = {
            "package": package,
            "available": importlib.util.find_spec(module) is not None,
        }
    missing = sorted(
        details["package"] for details in document_formats.values()
        if not details["available"]
    )
    return {
        "ocr_engines": engines,
        "ocr_languages": langs,
        "image_formats": sorted(_IMAGE_SUFFIXES),
        "text_formats": sorted(_TEXT_SUFFIXES),
        "document_formats": document_formats,
        "legacy_formats_requiring_conversion": dict(_LEGACY_FORMATS),
        "missing_libraries": missing,
        "confidence": "heuristic (word-shape/length/volume), not an engine score",
    }
